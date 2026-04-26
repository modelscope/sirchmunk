# Copyright (c) ModelScope Contributors. All rights reserved.
"""
Knowledge compiler — orchestrates offline compile of document collections.

Fuses PageIndex (tree indexing) and LLM Wiki (knowledge compilation network)
into a single compile pipeline that produces structured tree indices and
knowledge clusters for downstream search acceleration.
"""

import asyncio
import json
import math
import os
import random
import re
import hashlib
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple, Union

from sirchmunk.learnings.tree_indexer import (
    DocumentTree,
    DocumentTreeIndexer,
)
from sirchmunk.llm.openai_chat import OpenAIChat
from sirchmunk.schema.knowledge import (
    AbstractionLevel,
    EvidenceUnit,
    KnowledgeCluster,
    Lifecycle,
    WeakSemanticEdge,
)
from sirchmunk.storage.knowledge_storage import KnowledgeStorage
from sirchmunk.utils import LogCallback, create_logger
from sirchmunk.utils.document_extractor import DocumentExtractor
from sirchmunk.utils.file_utils import get_fast_hash

# Concurrency cap for LLM-heavy file processing
_DEFAULT_CONCURRENCY = 3

# Similarity threshold for merging into existing clusters during compile
_MERGE_SIMILARITY_THRESHOLD = 0.75

# Max chars for manifest-persisted document summary (used in Phase 2 & catalog)
_MANIFEST_SUMMARY_MAX_LEN = 500

# Preview window for direct LLM summarisation (no tree), ~4K tokens
_SUMMARY_PREVIEW_CHARS = 16_000

# Multi-section sampling for large documents without a tree index
_SUMMARY_SAMPLE_SECTIONS = 3          # Number of sections to sample for large docs
_SUMMARY_SAMPLE_SECTION_CHARS = 5_000  # Chars per sampled section

# Targeted table extraction: max chars per table region
_TARGETED_TABLE_MAX_CHARS = 5000

# Targeted table extraction: only process nodes spanning <= N pages
_TABLE_PAGE_SPAN_LIMIT = 5

# Numeric density threshold – fraction of numeric/symbol chars ($, %, digits,
# parenthesised numbers) relative to total non-whitespace chars.  Pages below
# this threshold are skipped during targeted extraction.
_TABLE_NUMERIC_DENSITY_THRESHOLD = 0.15

# Excel table-level adaptive sampling constants
_XLSX_TOTAL_ROW_BUDGET = 100       # Total sampled rows budget across all sheets
_XLSX_MIN_ROWS_PER_SHEET = 3       # Minimum sampled rows per sheet
_XLSX_MAX_ROWS_PER_SHEET = 50      # Maximum sampled rows per sheet
_XLSX_MAX_SHEETS = 10              # Maximum number of sheets to process
_XLSX_MAX_COLS_DISPLAY = 20        # Maximum columns to display per sheet


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class FileManifestEntry:
    """State of a single file in the compile manifest."""

    file_hash: str
    compiled_at: str
    has_tree: bool
    cluster_ids: List[str]
    size_bytes: int
    summary: str = ""  # 新增：存储编译期生成的文档摘要
    has_explicit_toc: bool = False  # Whether a native TOC was extracted from the file
    tree_node_count: int = 0  # Number of nodes in the tree index (quality metric)
    has_xlsx_digest: bool = False  # Whether a pre-compiled Excel evidence digest exists
    has_table_digest: bool = False  # Whether PDF tables were extracted and stored
    table_count: int = 0  # Number of tables in this file

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_hash": self.file_hash,
            "compiled_at": self.compiled_at,
            "has_tree": self.has_tree,
            "cluster_ids": self.cluster_ids,
            "size_bytes": self.size_bytes,
            "summary": self.summary,
            "has_explicit_toc": self.has_explicit_toc,
            "tree_node_count": self.tree_node_count,
            "has_xlsx_digest": self.has_xlsx_digest,
            "has_table_digest": self.has_table_digest,
            "table_count": self.table_count,
        }

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "FileManifestEntry":
        return cls(
            file_hash=data["file_hash"],
            compiled_at=data["compiled_at"],
            has_tree=data.get("has_tree", False),
            cluster_ids=data.get("cluster_ids", []),
            size_bytes=data.get("size_bytes", 0),
            summary=data.get("summary", ""),
            has_explicit_toc=data.get("has_explicit_toc", False),
            tree_node_count=data.get("tree_node_count", 0),
            has_xlsx_digest=data.get("has_xlsx_digest", False),
            has_table_digest=data.get("has_table_digest", False),
            table_count=data.get("table_count", 0),
        )


@dataclass
class CompileManifest:
    """Tracks compiled file states for incremental processing."""

    version: str = "1.0"
    last_compile_at: Optional[str] = None
    files: Dict[str, FileManifestEntry] = field(default_factory=dict)

    def to_json(self) -> str:
        return json.dumps({
            "version": self.version,
            "last_compile_at": self.last_compile_at,
            "files": {k: v.to_dict() for k, v in self.files.items()},
        }, ensure_ascii=False, indent=2)

    @classmethod
    def from_json(cls, json_str: str) -> "CompileManifest":
        data = json.loads(json_str)
        files = {
            k: FileManifestEntry.from_dict(v)
            for k, v in data.get("files", {}).items()
        }
        return cls(
            version=data.get("version", "1.0"),
            last_compile_at=data.get("last_compile_at"),
            files=files,
        )


@dataclass
class FileEntry:
    """Discovered file pending compilation."""

    path: str
    size_bytes: int
    file_hash: str


@dataclass
class ChangeSet:
    """Delta between discovered files and the manifest."""

    added: List[FileEntry] = field(default_factory=list)
    modified: List[FileEntry] = field(default_factory=list)
    deleted: List[str] = field(default_factory=list)
    unchanged: List[str] = field(default_factory=list)


@dataclass
class FileCompileResult:
    """Result of compiling a single file."""

    path: str
    tree: Optional[DocumentTree] = None
    summary: str = ""
    topics: List[str] = field(default_factory=list)
    evidence: Optional[EvidenceUnit] = None
    cluster_ids: List[str] = field(default_factory=list)
    error: Optional[str] = None
    has_explicit_toc: bool = False  # Whether TOC was extracted from native structure
    tree_node_count: int = 0  # Number of nodes in the tree index
    has_xlsx_digest: bool = False  # Whether a pre-compiled Excel evidence digest exists
    has_table_digest: bool = False  # Whether a pre-compiled table digest exists
    table_count: int = 0  # Number of tables extracted


@dataclass
class CompileReport:
    """Summary report of a compile run."""

    total_files: int = 0
    files_added: int = 0
    files_modified: int = 0
    files_skipped: int = 0
    files_deleted: int = 0
    files_sampled: int = 0
    trees_built: int = 0
    clusters_created: int = 0
    clusters_merged: int = 0
    cross_refs_built: int = 0
    errors: List[str] = field(default_factory=list)
    elapsed_seconds: float = 0.0

    def to_dict(self) -> Dict[str, Any]:
        return {
            "total_files": self.total_files,
            "files_added": self.files_added,
            "files_modified": self.files_modified,
            "files_skipped": self.files_skipped,
            "files_deleted": self.files_deleted,
            "files_sampled": self.files_sampled,
            "trees_built": self.trees_built,
            "clusters_created": self.clusters_created,
            "clusters_merged": self.clusters_merged,
            "cross_refs_built": self.cross_refs_built,
            "errors": self.errors,
            "elapsed_seconds": round(self.elapsed_seconds, 2),
        }


@dataclass
class CompileStatus:
    """Status snapshot of the compile state."""

    total_compiled_files: int = 0
    total_clusters: int = 0
    total_trees: int = 0
    last_compile_at: Optional[str] = None
    manifest_path: str = ""


# ---------------------------------------------------------------------------
# Importance probability sampler
# ---------------------------------------------------------------------------

class ImportanceSampler:
    """Select a representative subset of files using importance-based probability.

    Sampling strategy for large datasets:
    - Larger files get higher probability (they contain more information).
    - Uncompiled (new) files are prioritised over previously compiled ones.
    - Files with rare extensions get a mild boost (diversity signal).
    - The final probability is proportional to a composite importance score.
    """

    def __init__(self, max_files: int, seed: Optional[int] = None):
        self._max_files = max_files
        self._rng = random.Random(seed)

    def sample(self, files: List[FileEntry], manifest: CompileManifest) -> List[FileEntry]:
        """Return up to *max_files* entries sampled by importance."""
        if len(files) <= self._max_files:
            return files

        scores = [self._score(f, manifest) for f in files]
        total = sum(scores) or 1.0
        probs = [s / total for s in scores]

        selected_indices = set()
        attempts = 0
        while len(selected_indices) < self._max_files and attempts < len(files) * 3:
            idx = self._weighted_choice(probs)
            selected_indices.add(idx)
            attempts += 1

        return [files[i] for i in sorted(selected_indices)]

    def _score(self, entry: FileEntry, manifest: CompileManifest) -> float:
        """Compute composite importance score."""
        # Size factor: log-scaled, bounded
        size_score = math.log2(max(entry.size_bytes, 1024)) / 20.0

        # Novelty factor: new files are more important
        novelty = 2.0 if entry.path not in manifest.files else 0.5

        # Extension diversity: rare extensions get a mild boost
        ext = Path(entry.path).suffix.lower()
        diversity = 1.5 if ext in {".pdf", ".docx", ".doc", ".tex"} else 1.0

        return size_score * novelty * diversity

    def _weighted_choice(self, probs: List[float]) -> int:
        r = self._rng.random()
        cumulative = 0.0
        for i, p in enumerate(probs):
            cumulative += p
            if r <= cumulative:
                return i
        return len(probs) - 1


# ---------------------------------------------------------------------------
# Compiler
# ---------------------------------------------------------------------------

class KnowledgeCompiler:
    """Orchestrate compile pipeline: file discovery -> tree indexing -> knowledge aggregation."""

    # File extensions eligible for compilation
    _ELIGIBLE_EXTENSIONS = {
        ".pdf", ".docx", ".doc", ".md", ".markdown", ".html", ".htm",
        ".rst", ".tex", ".txt", ".pptx", ".xlsx",
    }

    def __init__(
        self,
        llm: OpenAIChat,
        embedding_client: Optional[Any],
        knowledge_storage: KnowledgeStorage,
        tree_indexer: DocumentTreeIndexer,
        work_path: Union[str, Path],
        log_callback: LogCallback = None,
    ):
        self._llm = llm
        self._embedding = embedding_client
        self._storage = knowledge_storage
        self._tree_indexer = tree_indexer
        self._work_path = Path(work_path).expanduser().resolve()
        self._log = create_logger(log_callback=log_callback)

        self._compile_dir = self._work_path / ".cache" / "compile"
        self._compile_dir.mkdir(parents=True, exist_ok=True)
        self._manifest_path = self._compile_dir / "manifest.json"

    # ------------------------------------------------------------------ #
    #  Public API                                                         #
    # ------------------------------------------------------------------ #

    async def compile(
        self,
        paths: List[str],
        *,
        incremental: bool = True,
        shallow: bool = False,
        max_files: Optional[int] = None,
        concurrency: int = _DEFAULT_CONCURRENCY,
    ) -> CompileReport:
        """Execute the unified knowledge compile pipeline.

        Args:
            paths: Directories or files to compile.
            incremental: Skip unchanged files.
            shallow: Skip tree building even for eligible files — use direct
                     LLM summarisation only (faster, lower quality).
            max_files: Cap on files to process (triggers importance sampling).
            concurrency: Max parallel file compilations.
        """
        import time
        t0 = time.monotonic()
        report = CompileReport()

        # Phase 1: discover and diff
        await self._log.info("[Compile] Phase 1: File discovery & change detection")
        manifest = self._load_manifest()
        discovered = await self._discover_files(paths)
        report.total_files = len(discovered)
        await self._log.info(f"[Compile] Discovered {len(discovered)} eligible files")

        if incremental:
            changes = self._detect_changes(discovered, manifest)
            to_compile = changes.added + changes.modified
            report.files_skipped = len(changes.unchanged)
            report.files_deleted = len(changes.deleted)
            for deleted_path in changes.deleted:
                manifest.files.pop(deleted_path, None)
        else:
            to_compile = discovered
            report.files_skipped = 0

        report.files_added = len([f for f in to_compile if f.path not in manifest.files])
        report.files_modified = len(to_compile) - report.files_added

        # Phase 1.5: importance sampling for large datasets
        if max_files and len(to_compile) > max_files:
            await self._log.info(
                f"[Compile] Applying importance sampling: {len(to_compile)} -> {max_files} files"
            )
            sampler = ImportanceSampler(max_files=max_files)
            to_compile = sampler.sample(to_compile, manifest)
            report.files_sampled = len(to_compile)

        if not to_compile:
            await self._log.info("[Compile] No files to compile (all up-to-date)")
            report.elapsed_seconds = time.monotonic() - t0
            return report

        await self._log.info(
            f"[Compile] Phase 2: Processing {len(to_compile)} files "
            f"(concurrency={concurrency})"
        )

        # Phase 2: compile files with bounded concurrency
        semaphore = asyncio.Semaphore(concurrency)
        results: List[FileCompileResult] = []

        async def _bounded(entry: FileEntry) -> FileCompileResult:
            async with semaphore:
                return await self._compile_single_file(entry, shallow=shallow)

        tasks = [_bounded(f) for f in to_compile]
        for coro in asyncio.as_completed(tasks):
            result = await coro
            results.append(result)
            if result.error:
                report.errors.append(f"{result.path}: {result.error}")
            else:
                if result.tree:
                    report.trees_built += 1
                # Update manifest
                manifest.files[result.path] = FileManifestEntry(
                    file_hash=get_fast_hash(result.path) or "",
                    compiled_at=datetime.now(timezone.utc).isoformat(),
                    has_tree=result.tree is not None,
                    cluster_ids=result.cluster_ids,
                    size_bytes=Path(result.path).stat().st_size if Path(result.path).exists() else 0,
                    summary=result.summary[:_MANIFEST_SUMMARY_MAX_LEN] if result.summary else "",
                    has_explicit_toc=result.has_explicit_toc,
                    tree_node_count=result.tree_node_count,
                    has_xlsx_digest=result.has_xlsx_digest,
                    has_table_digest=result.has_table_digest,
                    table_count=result.table_count,
                )
                _mentry = manifest.files[result.path]
                print(f"SEARCH_WIKI_DEBUG [C4] manifest_entry: has_tree={_mentry.has_tree}, has_table_digest={_mentry.has_table_digest}, file_hash={_mentry.file_hash}", flush=True)

        # Phase 3: aggregate results into knowledge network
        await self._log.info("[Compile] Phase 3: Knowledge aggregation")
        for r in results:
            if r.error or not r.summary:
                continue
            created, merged = await self._aggregate_to_knowledge_network(r)
            report.clusters_created += created
            report.clusters_merged += merged

        # Phase 4: cross-references
        await self._log.info("[Compile] Phase 4: Building cross-references")
        report.cross_refs_built = await self._build_cross_references(results)

        # Phase 5: persist manifest + document catalog
        manifest.last_compile_at = datetime.now(timezone.utc).isoformat()
        self._save_manifest(manifest)
        self._storage.force_sync()

        # Generate document catalog for search-time routing
        self._build_document_catalog(manifest)

        # Phase: Build summary index for embedding+BM25 fallback (optional, non-blocking)
        await self._build_summary_index(manifest)

        report.elapsed_seconds = time.monotonic() - t0
        await self._log.info(
            f"[Compile] Done in {report.elapsed_seconds:.1f}s — "
            f"trees={report.trees_built}, created={report.clusters_created}, "
            f"merged={report.clusters_merged}, errors={len(report.errors)}"
        )
        return report

    async def get_status(self, paths: List[str]) -> CompileStatus:
        """Return current compile status for the given paths."""
        manifest = self._load_manifest()
        path_set = {str(Path(p).resolve()) for p in paths}

        compiled_count = 0
        tree_count = 0
        cluster_ids: Set[str] = set()
        for fp, entry in manifest.files.items():
            for p in path_set:
                if fp.startswith(p):
                    compiled_count += 1
                    if entry.has_tree:
                        tree_count += 1
                    cluster_ids.update(entry.cluster_ids)
                    break

        return CompileStatus(
            total_compiled_files=compiled_count,
            total_clusters=len(cluster_ids),
            total_trees=tree_count,
            last_compile_at=manifest.last_compile_at,
            manifest_path=str(self._manifest_path),
        )

    # ------------------------------------------------------------------ #
    #  File discovery and change detection                                #
    # ------------------------------------------------------------------ #

    async def _discover_files(self, paths: List[str]) -> List[FileEntry]:
        """Walk paths and return all compilation-eligible files."""
        entries: List[FileEntry] = []
        seen: Set[str] = set()

        for base in paths:
            base_path = Path(base).expanduser().resolve()
            if base_path.is_file():
                candidates = [base_path]
            elif base_path.is_dir():
                candidates = sorted(base_path.rglob("*"))
            else:
                continue

            for fp in candidates:
                if not fp.is_file():
                    continue
                if fp.suffix.lower() not in self._ELIGIBLE_EXTENSIONS:
                    continue
                abs_path = str(fp.resolve())
                if abs_path in seen:
                    continue
                seen.add(abs_path)
                fh = get_fast_hash(abs_path)
                if fh is None:
                    continue
                entries.append(FileEntry(
                    path=abs_path,
                    size_bytes=fp.stat().st_size,
                    file_hash=fh,
                ))

        return entries

    def _detect_changes(
        self, discovered: List[FileEntry], manifest: CompileManifest,
    ) -> ChangeSet:
        """Compare discovered files against the manifest for incremental compile."""
        changes = ChangeSet()
        current_paths = {f.path for f in discovered}

        for entry in discovered:
            prev = manifest.files.get(entry.path)
            if prev is None:
                changes.added.append(entry)
            elif prev.file_hash != entry.file_hash:
                changes.modified.append(entry)
            else:
                changes.unchanged.append(entry.path)

        for old_path in manifest.files:
            if old_path not in current_paths:
                changes.deleted.append(old_path)

        return changes

    # ------------------------------------------------------------------ #
    #  Single-file compilation                                            #
    # ------------------------------------------------------------------ #

    async def _compile_single_file(
        self,
        entry: FileEntry,
        *,
        shallow: bool = False,
    ) -> FileCompileResult:
        """Unified compile pipeline: tree-if-eligible -> summary -> topics -> evidence.

        When *shallow* is True (or file is ineligible for tree indexing),
        the pipeline skips tree building and summarises via a direct LLM call.
        """
        result = FileCompileResult(path=entry.path)
        print(f"SEARCH_WIKI_DEBUG [C1] _compile_single_file: file_path={entry.path}, file_hash={entry.file_hash}", flush=True)
        try:
            await self._log.info(f"[Compile] Processing: {Path(entry.path).name}")

            extraction = await DocumentExtractor.extract(
                entry.path, DocumentExtractor.ENHANCED,
            )
            content = extraction.content
            if not content or len(content.strip()) < 100:
                result.error = "Insufficient text content"
                return result

            use_tree = (
                not shallow
                and DocumentTreeIndexer.should_build_tree(entry.path, len(content))
            )

            # Phase 0.5: TOC extraction (layers 1-3 are zero LLM calls)
            toc_entries = None
            if use_tree:
                from sirchmunk.learnings.toc_extractor import TOCExtractor
                toc_entries = await TOCExtractor.extract(
                    entry.path, content,
                    total_pages=extraction.page_count,
                )
                if toc_entries:
                    await self._log.info(
                        f"[Compile] Extracted TOC with {len(toc_entries)} entries "
                        f"for {Path(entry.path).name}"
                    )

            if use_tree:
                result.tree = await self._tree_indexer.build_tree(
                    entry.path, content,
                    toc_entries=toc_entries,
                    total_pages=extraction.page_count,
                )

            # Record TOC / tree metrics on the result for manifest persistence
            result.has_explicit_toc = toc_entries is not None and len(toc_entries) > 0
            result.tree_node_count = self._count_tree_nodes(result.tree)
            print(f"SEARCH_WIKI_DEBUG [C2] tree_build: success={result.tree is not None}, nodes={result.tree_node_count}, tree.file_path={result.tree.file_path if result.tree else 'N/A'}", flush=True)

            # Enrich content with structural metadata for non-text types
            ext = Path(entry.path).suffix.lower()
            evidence_digest = ""

            if ext in (".xlsx", ".xls"):
                # Excel: use adaptive sampling for both metadata and evidence
                metadata_prefix, evidence_digest = self._extract_xlsx_sampling(entry.path)
                enriched_content = metadata_prefix + content if metadata_prefix else content
            else:
                metadata_prefix = self._extract_structured_metadata(entry.path, content)
                enriched_content = metadata_prefix + content if metadata_prefix else content

            result.summary = await self._extract_summary(
                entry.path, enriched_content, result.tree,
            )
            result.topics = await self._extract_topics(result.summary)
            result.evidence = self._build_evidence(entry, content, result)

            # Persist Excel evidence digest for search-time consumption
            if evidence_digest.strip():
                try:
                    digest_dir = self._compile_dir / "xlsx_digests"
                    digest_dir.mkdir(parents=True, exist_ok=True)
                    file_hash = get_fast_hash(entry.path) or ""
                    if file_hash:
                        digest_path = digest_dir / f"{file_hash}.txt"
                        digest_path.write_text(evidence_digest, encoding="utf-8")
                        result.has_xlsx_digest = True
                except Exception:
                    pass

            # Persist table digest for documents with extracted tables
            if extraction.tables:
                try:
                    table_digest = self._build_table_digest(extraction.tables)
                    if table_digest:
                        digest_dir = self._compile_dir / "table_digests"
                        digest_dir.mkdir(parents=True, exist_ok=True)
                        file_hash = get_fast_hash(entry.path) or ""
                        if file_hash:
                            digest_path = digest_dir / f"{file_hash}.json"
                            digest_path.write_text(
                                json.dumps(table_digest, ensure_ascii=False),
                                encoding="utf-8",
                            )
                            result.has_table_digest = True
                            result.table_count = len(extraction.tables)
                except Exception:
                    pass

            print(f"SEARCH_WIKI_DEBUG [C3] table_digest: generated={result.has_table_digest}, count={result.table_count}", flush=True)

            # Integrate tables into tree: annotate counts + create table child nodes
            if result.tree and result.tree.root and extraction.tables:
                self._integrate_tables_into_tree(
                    result.tree.root, extraction.tables,
                    content=content, total_pages=extraction.page_count,
                )

            # Phase 2.5: Targeted table extraction via generic structural signals
            if result.tree and result.tree.root and ext == ".pdf":
                targeted_tables = await self._targeted_table_extraction(
                    entry.path, result.tree,
                )
                if targeted_tables:
                    # Load existing table digest (if any) and merge
                    digest_dir = self._compile_dir / "table_digests"
                    file_hash = get_fast_hash(entry.path) or ""
                    existing_digest: list[dict] = []
                    if file_hash and result.has_table_digest:
                        digest_path = digest_dir / f"{file_hash}.json"
                        if digest_path.exists():
                            try:
                                raw = json.loads(
                                    digest_path.read_text(encoding="utf-8")
                                )
                                existing_digest = raw.get("tables", [])
                            except Exception:
                                pass
                    merged = self._merge_table_digests(
                        existing_digest, targeted_tables,
                    )
                    if merged and file_hash:
                        digest_dir.mkdir(parents=True, exist_ok=True)
                        digest_path = digest_dir / f"{file_hash}.json"
                        digest_path.write_text(
                            json.dumps(
                                {
                                    "version": 1,
                                    "table_count": len(merged),
                                    "tables": merged,
                                },
                                ensure_ascii=False,
                            ),
                            encoding="utf-8",
                        )
                        result.has_table_digest = True
                        result.table_count = len(merged)
                        await self._log.info(
                            f"[Compile] Targeted table extraction added "
                            f"{len(targeted_tables)} tables for "
                            f"{Path(entry.path).name}"
                        )

        except Exception as exc:
            result.error = str(exc)
            await self._log.warning(f"[Compile] Failed: {entry.path}: {exc}")

        return result

    async def _extract_summary(
        self,
        file_path: str,
        content: str,
        tree: Optional[DocumentTree] = None,
    ) -> str:
        """Generate a document-level summary.

        When a tree is available its root already contains an LLM-synthesized
        summary (produced by ``_synthesize_root_summary`` during tree build),
        so we reuse it directly — no redundant LLM call.

        For large documents without a tree, uses multi-section sampling
        (beginning, middle, end) to capture the full scope of the document.
        """
        if tree and tree.root and tree.root.summary:
            return tree.root.summary

        preview = self._build_summary_preview(content)
        from sirchmunk.llm.prompts import COMPILE_DOC_SUMMARY
        prompt = COMPILE_DOC_SUMMARY.format(
            file_name=Path(file_path).name,
            document_content=preview,
        )
        resp = await self._llm.achat([{"role": "user", "content": prompt}])
        return resp.content.strip()

    @staticmethod
    def _build_summary_preview(content: str) -> str:
        """Build a representative preview for LLM summarisation.

        For short documents (≤ _SUMMARY_PREVIEW_CHARS), returns the full
        content.  For large documents, samples the beginning, middle, and
        end to capture the document's full scope within the token budget.
        """
        if len(content) <= _SUMMARY_PREVIEW_CHARS:
            return content

        section_size = _SUMMARY_SAMPLE_SECTION_CHARS
        mid_start = max(section_size, (len(content) - section_size) // 2)

        head = content[:section_size]
        middle = content[mid_start:mid_start + section_size]
        tail = content[-section_size:]

        return (
            f"[Beginning of document]\n{head}\n\n"
            f"[... content omitted ...]\n\n"
            f"[Middle of document]\n{middle}\n\n"
            f"[... content omitted ...]\n\n"
            f"[End of document]\n{tail}"
        )

    @staticmethod
    def _extract_structured_metadata(file_path: str, content: str) -> str:
        """Extract structural metadata for non-text document types.

        For spreadsheets and presentations, prepend a structural overview
        (sheet names, column headers, slide titles) so the LLM summariser
        has better context than raw extracted text alone.

        Returns a metadata prefix string (may be empty for unsupported types).
        """
        ext = Path(file_path).suffix.lower()

        if ext == ".xlsx":
            metadata, _evidence = KnowledgeCompiler._extract_xlsx_sampling(file_path)
            return metadata
        if ext == ".pptx":
            return KnowledgeCompiler._extract_pptx_metadata(file_path)

        return ""

    @staticmethod
    def _compute_xlsx_sample_rows(total_rows: int, num_sheets: int, sheet_rows: int) -> int:
        """Compute adaptive sample row count per sheet.

        Strategy:
        - Divides _XLSX_TOTAL_ROW_BUDGET equally across sheets
        - Small sheets (<=budget) are fully sampled
        - Large sheets are capped at budget
        - Result clamped to [_XLSX_MIN_ROWS_PER_SHEET, _XLSX_MAX_ROWS_PER_SHEET]
        """
        budget_per_sheet = max(1, _XLSX_TOTAL_ROW_BUDGET // max(1, num_sheets))
        n = min(sheet_rows, budget_per_sheet)
        return max(_XLSX_MIN_ROWS_PER_SHEET, min(_XLSX_MAX_ROWS_PER_SHEET, n))

    @staticmethod
    def _extract_xlsx_sampling(file_path: str) -> Tuple[str, str]:
        """Extract structural metadata AND sampled content from Excel workbook.

        Performs table-level intelligent sampling with adaptive row counts
        based on workbook size and sheet complexity.

        Returns:
            (metadata_prefix, evidence_digest)
            - metadata_prefix: injected into summary generation context
            - evidence_digest: structured text usable directly as search evidence
        """
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

            sheet_names = wb.sheetnames[:_XLSX_MAX_SHEETS]
            num_sheets = len(sheet_names)

            # Phase 1: Collect sheet statistics
            sheet_stats: List[Dict[str, Any]] = []
            for sheet_name in sheet_names:
                ws = wb[sheet_name]
                row_count = ws.max_row or 0
                col_count = ws.max_column or 0
                # Read headers (first row)
                headers: List[str] = []
                for row in ws.iter_rows(min_row=1, max_row=1, values_only=True):
                    headers = [str(h) for h in row if h is not None]
                    break
                sheet_stats.append({
                    "name": sheet_name,
                    "rows": row_count,
                    "cols": col_count,
                    "headers": headers[:_XLSX_MAX_COLS_DISPLAY],
                    "ws": ws,
                })

            # Phase 2: Calculate total rows for adaptive sampling
            total_rows = sum(s["rows"] for s in sheet_stats)

            meta_lines: List[str] = ["[Excel Workbook Structure]"]
            evidence_lines: List[str] = []

            for stat in sheet_stats:
                ws = stat["ws"]
                sheet_name = stat["name"]
                row_count = stat["rows"]
                col_count = stat["cols"]
                headers = stat["headers"]
                header_str = ", ".join(headers) if headers else "no headers"

                # Metadata line
                meta_lines.append(
                    f"- Sheet '{sheet_name}': {row_count} rows, {col_count} columns, "
                    f"headers: [{header_str}]"
                )

                # Adaptive sampling
                sample_n = KnowledgeCompiler._compute_xlsx_sample_rows(
                    total_rows, num_sheets, row_count
                )

                evidence_lines.append(
                    f"[Sheet '{sheet_name}' ({row_count} rows, {col_count} columns)]"
                )
                evidence_lines.append(f"Columns: {header_str}")

                # Sample rows
                if row_count <= sample_n:
                    evidence_lines.append(f"(Full content - {row_count} rows)")
                else:
                    evidence_lines.append(f"Sample rows (top {sample_n} of {row_count}):")

                # Build table header
                display_headers = headers[:_XLSX_MAX_COLS_DISPLAY]
                if display_headers:
                    evidence_lines.append("| " + " | ".join(display_headers) + " |")
                    evidence_lines.append("|" + "|".join(["---"] * len(display_headers)) + "|")

                # Read sample rows (skip header row)
                numeric_cols: Dict[int, List[float]] = {}  # col_index -> numeric values
                sampled = 0
                for row in ws.iter_rows(
                    min_row=2,
                    max_row=min(row_count, sample_n + 1),
                    values_only=True,
                ):
                    cells: List[str] = []
                    for ci, cell_val in enumerate(row):
                        if ci >= _XLSX_MAX_COLS_DISPLAY:
                            break
                        str_val = str(cell_val) if cell_val is not None else ""
                        cells.append(str_val[:50])  # truncate long cell values
                        # Track numeric values for statistics
                        if isinstance(cell_val, (int, float)) and cell_val == cell_val:
                            numeric_cols.setdefault(ci, []).append(float(cell_val))
                    if cells:
                        evidence_lines.append("| " + " | ".join(cells) + " |")
                    sampled += 1

                # Statistics for numeric columns
                stat_parts: List[str] = []
                for ci, values in numeric_cols.items():
                    if len(values) >= 2 and ci < len(display_headers):
                        col_name = display_headers[ci]
                        stat_parts.append(
                            f"{col_name} range [{min(values):.4g}-{max(values):.4g}]"
                        )
                if stat_parts:
                    evidence_lines.append(f"Statistics: {', '.join(stat_parts[:5])}")

                evidence_lines.append("")  # blank line between sheets

            wb.close()

            metadata = "\n".join(meta_lines) + "\n\n"
            evidence = "\n".join(evidence_lines)
            return metadata, evidence

        except Exception:
            return "", ""

    @staticmethod
    def _extract_xlsx_metadata(file_path: str) -> str:
        """Extract structural metadata from Excel files (legacy wrapper).

        Delegates to _extract_xlsx_sampling and returns only the metadata prefix
        for backward compatibility.
        """
        metadata, _evidence = KnowledgeCompiler._extract_xlsx_sampling(file_path)
        return metadata

    @staticmethod
    def _extract_pptx_metadata(file_path: str) -> str:
        """Extract structural metadata from PowerPoint files.

        Reads slide count and titles (from the title placeholder) to give
        the LLM a table-of-contents-like overview of the presentation.
        Caps at 20 slides for bounded output.
        """
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            lines: List[str] = [f"[PowerPoint Structure: {len(prs.slides)} slides]"]
            for i, slide in enumerate(prs.slides[:20], 1):  # Cap at 20 slides
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                if title:
                    lines.append(f"- Slide {i}: {title}")
            return "\n".join(lines) + "\n\n"
        except Exception:
            return ""

    def _build_evidence(
        self,
        entry: FileEntry,
        content: str,
        result: FileCompileResult,
    ) -> EvidenceUnit:
        """Build an EvidenceUnit, populating snippets/tree_path from tree leaves."""
        from sirchmunk.schema.metadata import FileInfo

        snippets: List[str] = []
        tree_path: Optional[List[str]] = None

        if result.tree and result.tree.root:
            leaves = result.tree.root.all_leaves()
            tree_path = [leaf.node_id for leaf in leaves]
            for leaf in leaves:
                start, end = leaf.char_range
                snippet = content[start:end][:500]
                if snippet.strip():
                    snippets.append(snippet)

        return EvidenceUnit(
            doc_id=FileInfo.get_cache_key(entry.path),
            file_or_url=Path(entry.path),
            summary=result.summary,
            is_found=True,
            snippets=snippets,
            tree_path=tree_path,
            extracted_at=datetime.now(timezone.utc),
        )

    async def _extract_topics(self, summary: str) -> List[str]:
        """Extract key topics/entities from a document summary."""
        from sirchmunk.llm.prompts import COMPILE_TOPIC_EXTRACTION
        prompt = COMPILE_TOPIC_EXTRACTION.format(summary=summary)
        resp = await self._llm.achat([{"role": "user", "content": prompt}])
        try:
            raw = resp.content.strip()
            if raw.startswith("["):
                parsed = json.loads(raw)
                if isinstance(parsed, list):
                    return [str(t) for t in parsed if t]
            return [t.strip() for t in raw.split(",") if t.strip()]
        except (json.JSONDecodeError, TypeError):
            return []

    # ------------------------------------------------------------------ #
    #  Knowledge aggregation (LLM Wiki Ingest)                            #
    # ------------------------------------------------------------------ #

    async def _aggregate_to_knowledge_network(
        self, result: FileCompileResult,
    ) -> Tuple[int, int]:
        """Aggregate a file's compile result into the knowledge network.

        Three-tier similarity strategy (per design doc):
          - similarity >= 0.80  → merge into existing cluster
          - 0.50 <= sim < 0.80  → create new cluster + weak edge to similar
          - similarity < 0.50   → create standalone cluster

        Returns:
            (clusters_created, clusters_merged)
        """
        created, merged = 0, 0
        if not result.summary:
            return created, merged

        embedding = self._encode_text(result.summary)

        # Search for similar existing clusters across a wider range
        best_match: Optional[Dict[str, Any]] = None
        if embedding is not None:
            similar = await self._storage.search_similar_clusters(
                query_embedding=embedding,
                top_k=3,
                similarity_threshold=0.50,
            )
            if similar:
                best_match = similar[0]

        if best_match and best_match["similarity"] >= 0.80:
            # Tier 1: merge into existing cluster
            cluster = await self._storage.get(best_match["id"])
            if cluster:
                await self._merge_into_cluster(cluster, result)
                # Re-compute embedding for merged content
                await self._update_cluster_embedding(cluster)
                result.cluster_ids.append(cluster.id)
                merged += 1
                return created, merged

        # Create a new cluster (Tier 2 or Tier 3)
        cluster = await self._create_cluster(result)
        if cluster:
            result.cluster_ids.append(cluster.id)
            await self._store_cluster_embedding(cluster, embedding, result.summary)
            created += 1

            # Tier 2: build weak edges to moderately similar clusters
            if best_match and best_match["similarity"] >= 0.50:
                for s in (similar or []):
                    if s["similarity"] >= 0.50:
                        target = await self._storage.get(s["id"])
                        if target:
                            self._add_edge(cluster, target.id, "embed_sim", s["similarity"])
                            self._add_edge(target, cluster.id, "embed_sim", s["similarity"])
                            await self._storage.update(target)
                await self._storage.update(cluster)

        return created, merged

    def _encode_text(self, text: str) -> Optional[Any]:
        """Encode text to embedding vector, returns None on failure."""
        if not self._embedding:
            return None
        try:
            return self._embedding.encode(text)
        except Exception:
            return None

    async def _store_cluster_embedding(
        self, cluster: KnowledgeCluster, embedding: Optional[Any], text: str,
    ) -> None:
        """Store embedding for a cluster if available."""
        if embedding is None or not self._embedding:
            return
        text_hash = hashlib.md5(text.encode()).hexdigest()
        vec = embedding.tolist() if hasattr(embedding, "tolist") else list(embedding)
        await self._storage.store_embedding(
            cluster.id, vec,
            self._embedding.model_id or "default",
            text_hash,
        )

    async def _update_cluster_embedding(self, cluster: KnowledgeCluster) -> None:
        """Re-compute and store embedding after content merge."""
        content_text = str(cluster.content)[:2000] if cluster.content else ""
        if not content_text:
            return
        embedding = self._encode_text(content_text)
        await self._store_cluster_embedding(cluster, embedding, content_text)

    async def _merge_into_cluster(
        self,
        cluster: KnowledgeCluster,
        result: FileCompileResult,
    ) -> None:
        """Merge a file compile result into an existing cluster."""
        # Append evidence
        if result.evidence:
            existing_doc_ids = {e.doc_id for e in cluster.evidences}
            if result.evidence.doc_id not in existing_doc_ids:
                cluster.evidences.append(result.evidence)

        # Enrich content via LLM merge
        from sirchmunk.llm.prompts import COMPILE_MERGE_KNOWLEDGE
        prompt = COMPILE_MERGE_KNOWLEDGE.format(
            existing_content=str(cluster.content)[:3000],
            new_summary=result.summary[:3000],
        )
        resp = await self._llm.achat([{"role": "user", "content": prompt}])
        cluster.content = resp.content.strip()

        # Update metadata
        cluster.search_results = list(set(
            (cluster.search_results or []) + [result.path]
        ))
        merge_count = getattr(cluster, "merge_count", 0) or 0
        cluster.merge_count = merge_count + 1

        # Lifecycle promotion
        if cluster.merge_count >= 3 and cluster.lifecycle == Lifecycle.EMERGING:
            cluster.lifecycle = Lifecycle.STABLE

        await self._storage.update(cluster)

    async def _create_cluster(
        self, result: FileCompileResult,
    ) -> Optional[KnowledgeCluster]:
        """Create a new KnowledgeCluster from a file compile result."""
        cluster_text = result.summary
        cluster_id = f"C{hashlib.sha256(cluster_text.encode('utf-8')).hexdigest()[:10]}"

        name = Path(result.path).stem[:60]
        if result.topics:
            name = result.topics[0][:60]

        cluster = KnowledgeCluster(
            id=cluster_id,
            name=name,
            description=[result.summary[:500]],
            content=result.summary,
            evidences=[result.evidence] if result.evidence else [],
            patterns=result.topics[:5],
            lifecycle=Lifecycle.EMERGING,
            confidence=0.5,
            abstraction_level=AbstractionLevel.TECHNIQUE,
            hotness=0.3,
            search_results=[result.path],
        )

        ok = await self._storage.insert(cluster)
        return cluster if ok else None

    # ------------------------------------------------------------------ #
    #  Cross-references                                                   #
    # ------------------------------------------------------------------ #

    async def _build_cross_references(
        self, results: List[FileCompileResult],
    ) -> int:
        """Build co-occurrence edges between clusters that share source files.

        Two clusters are co-occurring when the same source file contributed
        evidence to both (e.g., different sections compiled into different
        clusters).  Includes historical data from the manifest.
        """
        # Build a complete map: cluster_id -> set of source file paths
        cluster_to_files: Dict[str, Set[str]] = {}

        # From current compile results
        for r in results:
            for cid in r.cluster_ids:
                cluster_to_files.setdefault(cid, set()).add(r.path)

        # From manifest (historical data)
        manifest = self._load_manifest()
        for fp, entry in manifest.files.items():
            for cid in entry.cluster_ids:
                cluster_to_files.setdefault(cid, set()).add(fp)

        # Find cluster pairs that share at least one source file
        cluster_ids = list(cluster_to_files.keys())
        edges_created = 0
        pairs_seen: Set[Tuple[str, str]] = set()

        for i in range(len(cluster_ids)):
            for j in range(i + 1, len(cluster_ids)):
                cid_a, cid_b = cluster_ids[i], cluster_ids[j]
                shared = cluster_to_files[cid_a] & cluster_to_files[cid_b]
                if not shared:
                    continue

                pair_key = (min(cid_a, cid_b), max(cid_a, cid_b))
                if pair_key in pairs_seen:
                    continue
                pairs_seen.add(pair_key)

                weight = min(len(shared) * 0.25, 1.0)
                c_a = await self._storage.get(cid_a)
                c_b = await self._storage.get(cid_b)
                if c_a and c_b:
                    self._add_edge(c_a, cid_b, "co_occur", weight)
                    self._add_edge(c_b, cid_a, "co_occur", weight)
                    await self._storage.update(c_a)
                    await self._storage.update(c_b)
                    edges_created += 1

        return edges_created

    @staticmethod
    def _add_edge(
        cluster: KnowledgeCluster, target_id: str, source: str, weight: float,
    ) -> None:
        """Add or update a WeakSemanticEdge on a cluster."""
        for edge in cluster.related_clusters:
            if edge.target_cluster_id == target_id and edge.source == source:
                edge.weight = max(edge.weight, weight)
                return
        cluster.related_clusters.append(
            WeakSemanticEdge(target_cluster_id=target_id, weight=weight, source=source)
        )

    def _build_table_digest(
        self, tables: List[Dict[str, Any]],
    ) -> Optional[Dict[str, Any]]:
        """Build a structured table digest from extraction output.

        Returns a versioned JSON-serializable dict containing all tables
        with their page numbers, markdown representation, and cell data.
        Tables are indexed for page-range-based retrieval at search time.
        """
        if not tables:
            return None

        digest_tables = []
        for idx, table in enumerate(tables):
            markdown = table.get("markdown", "")
            cells = table.get("cells", [])
            if not markdown and not cells:
                continue

            # Compute row/col counts from cells (kreuzberg returns List[List[str]])
            row_count = 0
            col_count = 0
            if cells:
                row_count = len(cells)
                col_count = max((len(row) for row in cells if isinstance(row, (list, tuple))), default=0)
            elif markdown:
                # Estimate from markdown lines
                lines = [l for l in markdown.strip().split("\n") if l.strip().startswith("|")]
                row_count = max(0, len(lines) - 1)  # exclude separator
                col_count = lines[0].count("|") - 1 if lines else 0

            # Skip pseudo-tables: single-column or insufficient structure
            if col_count <= 1:
                continue

            digest_tables.append({
                "index": idx,
                "page_number": table.get("page_number"),
                "markdown": markdown,
                "row_count": row_count,
                "col_count": col_count,
                "cells": cells,
            })

        if not digest_tables:
            return None

        return {
            "version": 1,
            "table_count": len(digest_tables),
            "tables": digest_tables,
        }

    def _integrate_tables_into_tree(
        self,
        node: "TreeNode",
        tables: List[Dict[str, Any]],
        content: str,
        *,
        total_pages: Optional[int] = None,
        _counter: Optional[List[int]] = None,
    ) -> None:
        """Integrate tables into tree: annotate counts AND create table child nodes for leaf nodes.

        For each node with a valid page_range, counts how many valid extracted
        tables fall within that range (excluding pseudo-tables with col_count <= 1).
        For leaf nodes with matching tables, creates dedicated TreeNode children
        with ``content_type="table"``.
        """
        from sirchmunk.learnings.tree_indexer import TreeNode

        if node is None:
            return

        if _counter is None:
            _counter = [0]

        # Depth-first: process existing children first
        for child in list(node.children):
            self._integrate_tables_into_tree(
                child, tables, content,
                total_pages=total_pages, _counter=_counter,
            )

        # Match valid tables to this node's page_range
        matched_tables: List[Dict[str, Any]] = []
        if node.page_range:
            ps, pe = node.page_range
            for t in tables:
                pn = t.get("page_number")
                if pn is None or not (ps <= pn <= pe):
                    continue
                # Skip pseudo-tables
                if self._is_pseudo_table(t):
                    continue
                matched_tables.append(t)

        node.table_count = len(matched_tables)

        # NOTE: _spawn_table_children disabled - converting leaf to non-leaf breaks
        # search navigation which expects leaves for char_range extraction.
        # TODO: Re-enable when search can properly handle mixed text+table children.
        # if not node.children and matched_tables:
        #     try:
        #         self._spawn_table_children(
        #             node, matched_tables, content, _counter,
        #         )
        #     except Exception:
        #         pass

    @staticmethod
    def _is_pseudo_table(table: Dict[str, Any]) -> bool:
        """Return True if the table lacks meaningful structure (col_count <= 1)."""
        markdown = table.get("markdown", "")
        cells = table.get("cells", [])
        if not markdown and not cells:
            return True
        col_count = 0
        if cells:
            col_count = max(
                (len(row) for row in cells if isinstance(row, (list, tuple))),
                default=0,
            )
        elif markdown:
            lines = [l for l in markdown.strip().split("\n") if l.strip().startswith("|")]
            col_count = (lines[0].count("|") - 1) if lines else 0
        return col_count <= 1

    def _spawn_table_children(
        self,
        node: "TreeNode",
        matched_tables: List[Dict[str, Any]],
        content: str,
        counter: List[int],
    ) -> None:
        """Create TreeNode children for each matched table under a leaf node.

        Also inserts a text-content sibling preserving the original leaf content.
        """
        from sirchmunk.learnings.tree_indexer import TreeNode

        child_level = node.level + 1

        # Preserve original text content as first child
        text_child_id = f"T{counter[0]:06d}"
        counter[0] += 1
        node.children.append(
            TreeNode(
                node_id=text_child_id,
                title=node.title,
                summary=node.summary[:300] if node.summary else "",
                char_range=node.char_range,
                level=child_level,
                page_range=node.page_range,
                children=[],
                table_count=0,
                content_type="text",
            )
        )

        # Create one child per table
        for table in matched_tables:
            tid = f"T{counter[0]:06d}"
            counter[0] += 1

            markdown = table.get("markdown", "")
            title = self._extract_table_title(table)
            page_number = table.get("page_number")

            # Attempt to locate table markdown in content
            char_range = node.char_range
            if markdown and content:
                pos = content.find(markdown[:120])
                if pos >= 0:
                    char_range = (pos, pos + len(markdown))

            page_range = (
                (page_number, page_number) if page_number is not None
                else node.page_range
            )

            node.children.append(
                TreeNode(
                    node_id=tid,
                    title=title,
                    summary=markdown[:300] if markdown else "",
                    char_range=char_range,
                    level=child_level,
                    page_range=page_range,
                    children=[],
                    table_count=0,
                    content_type="table",
                )
            )

    @staticmethod
    def _extract_table_title(table: Dict[str, Any]) -> str:
        """Extract a concise title from table markdown header row.

        Parses the first meaningful line of the markdown table (skipping
        separator rows like ``|---|---|``), strips ``|`` delimiters, and
        returns the first 80 characters as the title.
        """
        markdown = table.get("markdown", "")
        if not markdown:
            pn = table.get("page_number", "?")
            return f"Table (p.{pn})"

        for line in markdown.strip().split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            # Skip separator rows (e.g. |---|---| or +---+---+)
            content_chars = stripped.replace("|", "").replace("-", "").replace(":", "").replace("+", "").strip()
            if not content_chars:
                continue
            # Extract cell contents
            title = " | ".join(
                seg.strip() for seg in stripped.split("|") if seg.strip()
            )
            return title[:80] if title else f"Table (p.{table.get('page_number', '?')})"

        pn = table.get("page_number", "?")
        return f"Table (p.{pn})"

    @staticmethod
    def _count_tree_nodes(tree: Optional[DocumentTree]) -> int:
        """Count total nodes in a DocumentTree (recursive).

        Args:
            tree: The tree to count, or None.

        Returns:
            Total node count, or 0 if tree is None.
        """
        if tree is None or tree.root is None:
            return 0

        def _count(node: Any) -> int:
            return 1 + sum(_count(c) for c in node.children)

        return _count(tree.root)

    # ------------------------------------------------------------------ #
    #  Targeted table extraction                                          #
    # ------------------------------------------------------------------ #

    async def _targeted_table_extraction(
        self, file_path: str, tree: DocumentTree,
    ) -> list[dict]:
        """Extract tables from tree nodes likely containing tabular data.

        Uses generic structural signals (metadata, page span, numeric
        density) instead of domain-specific title keywords.  For each
        candidate with a valid ``page_range``, extracts per-page text
        via :meth:`DocumentExtractor.extract_page_range` and applies
        heuristic table-region detection.  Pages whose numeric density
        falls below ``_TABLE_NUMERIC_DENSITY_THRESHOLD`` are skipped.

        Returns:
            List of table dicts compatible with the table-digest format::

                {"page": int, "content": str, "source": str}
        """
        if tree is None or tree.root is None:
            return []

        candidates = self._find_table_candidate_nodes(tree.root)
        if not candidates:
            return []

        await self._log.info(
            f"[Compile] Targeted extraction: {len(candidates)} candidate "
            f"nodes in {Path(file_path).name}"
        )

        results: list[dict] = []
        seen_pages: set[int] = set()

        for node in candidates:
            if node.page_range is None:
                continue
            start_page, end_page = node.page_range
            # Skip pages already processed by another candidate
            page_nums = [p for p in range(start_page, end_page + 1)
                         if p not in seen_pages]
            if not page_nums:
                continue

            try:
                pages = DocumentExtractor.extract_page_range(
                    file_path, start_page, end_page,
                )
            except Exception as exc:
                await self._log.warning(
                    f"[Compile] Targeted extraction page read failed "
                    f"({start_page}-{end_page}): {exc}"
                )
                continue

            for pc in pages:
                if pc.page_number in seen_pages:
                    continue
                seen_pages.add(pc.page_number)
                # Numeric density gate – skip pages unlikely to contain tables
                if not self._page_has_table_density(pc.content):
                    continue
                regions = self._identify_table_regions(pc.content)
                for region in regions:
                    truncated = region[:_TARGETED_TABLE_MAX_CHARS]
                    results.append({
                        "page": pc.page_number,
                        "content": truncated,
                        "source": f"targeted:{node.title[:80]}",
                    })

        return results

    def _find_table_candidate_nodes(
        self, root: "TreeNode",
    ) -> list["TreeNode"]:
        """Collect leaf nodes that likely contain tables.

        Uses generic, domain-agnostic structural signals (any match
        suffices):

        - ``node.content_type == "table"`` – already tagged during compile.
        - ``node.table_count > 0`` – known to contain tables.
        - Has a valid ``page_range`` with span ≤ ``_TABLE_PAGE_SPAN_LIMIT``.
        """
        candidates: list = []

        def _walk(node: "TreeNode") -> None:
            if node.leaf:
                # Signal 1: content_type marked as table
                if getattr(node, "content_type", None) == "table":
                    candidates.append(node)
                    return
                # Signal 2: known to contain tables
                if getattr(node, "table_count", 0) > 0:
                    candidates.append(node)
                    return
                # Signal 3: moderate page span (tables rarely span many pages)
                page_range = getattr(node, "page_range", None)
                if page_range and len(page_range) == 2:
                    span = page_range[1] - page_range[0] + 1
                    if 1 <= span <= _TABLE_PAGE_SPAN_LIMIT:
                        candidates.append(node)
            else:
                for child in node.children:
                    _walk(child)

        _walk(root)
        return candidates

    @staticmethod
    def _page_has_table_density(page_text: str) -> bool:
        """Return True if *page_text* has numeric density above the threshold.

        Counts digits and common table symbols (``$``, ``%``, ``(``, ``)``)
        relative to total non-whitespace characters.
        """
        if not page_text:
            return False
        non_ws = sum(1 for ch in page_text if not ch.isspace())
        if non_ws == 0:
            return False
        numeric_chars = sum(
            1 for ch in page_text
            if ch.isdigit() or ch in "$%(),.+-"
        )
        return (numeric_chars / non_ws) >= _TABLE_NUMERIC_DENSITY_THRESHOLD

    @staticmethod
    def _identify_table_regions(page_text: str) -> list[str]:
        """Identify contiguous table-like regions in *page_text*.

        Heuristic rules:
        - Lines containing multiple numeric tokens (dollar amounts, %,
          parenthesised negatives) are considered *numeric rows*.
        - A run of >= 3 consecutive numeric rows forms a table region.
        - Leading/trailing whitespace rows are trimmed.

        Returns:
            List of extracted region strings (may be empty).
        """
        if not page_text:
            return []

        # Pattern: line has at least 2 numeric-looking tokens
        _NUM_TOKEN = re.compile(
            r"(?:"
            r"[\$€£¥]\s*[\d,.]+|"
            r"\([\d,.]+\)|"
            r"[\d,.]+%|"
            r"[\d]+\.[\d]+(?:[eE][+-]?\d+)?|"
            r"[\d,]{2,}"
            r")"
        )
        _MIN_NUMS_PER_LINE = 2
        _MIN_CONSECUTIVE = 3

        lines = page_text.split("\n")
        is_numeric = [
            len(_NUM_TOKEN.findall(line)) >= _MIN_NUMS_PER_LINE
            for line in lines
        ]

        regions: list[str] = []
        run_start: int | None = None

        for i, flag in enumerate(is_numeric):
            if flag:
                if run_start is None:
                    run_start = i
            else:
                if run_start is not None:
                    run_len = i - run_start
                    if run_len >= _MIN_CONSECUTIVE:
                        # Include one context line above/below
                        start = max(0, run_start - 1)
                        end = min(len(lines), i + 1)
                        regions.append(
                            "\n".join(lines[start:end]).strip()
                        )
                    run_start = None

        # Flush trailing run
        if run_start is not None:
            run_len = len(lines) - run_start
            if run_len >= _MIN_CONSECUTIVE:
                start = max(0, run_start - 1)
                regions.append(
                    "\n".join(lines[start:]).strip()
                )

        return regions

    @staticmethod
    def _get_table_page(entry: dict) -> int | None:
        """统一获取表格条目的页码，兼容 page_number 和 page 两种字段名。"""
        p = entry.get("page_number") or entry.get("page")
        return int(p) if p is not None else None

    @classmethod
    def _merge_table_digests(
        cls, existing: list[dict], new_tables: list[dict],
    ) -> list[dict]:
        """Merge *new_tables* into *existing* digest, deduplicating by page.

        If an existing entry and a new entry share the same page number,
        the new entry is skipped (existing kreuzberg-detected table takes
        precedence because it has richer structure like cells/markdown).

        Returns:
            Merged list suitable for storage in the table-digest JSON.
        """
        existing_pages = {cls._get_table_page(e) for e in existing}
        existing_pages.discard(None)

        merged = list(existing)
        for tbl in new_tables:
            page = cls._get_table_page(tbl)
            if page is not None and page in existing_pages:
                continue
            # Normalise to digest table format for consistency
            merged.append({
                "page_number": page,
                "markdown": tbl.get("content", ""),
                "row_count": None,
                "col_count": None,
                "cells": [],
                "source": tbl.get("source", "targeted"),
            })
        return merged

    # ------------------------------------------------------------------ #
    #  Summary index for embedding + BM25 fallback                        #
    # ------------------------------------------------------------------ #

    async def _build_summary_index(self, manifest: CompileManifest) -> None:
        """Build summary embedding + BM25 index for fallback search.

        Creates a lightweight index mapping each compiled file to:
        - Its summary text
        - Pre-computed embedding vector (384-dim, if EmbeddingUtil available)
        - Tokenized summary with term frequencies (via TokenizerUtil)

        The index is saved to .cache/compile/summary_index.json and consumed
        by search.py as a last-resort fallback when rga keyword search fails.

        Skips gracefully if dependencies (EmbeddingUtil/TokenizerUtil) are unavailable.
        """
        try:
            from sirchmunk.utils.tokenizer_util import TokenizerUtil
            from sirchmunk.learnings.summary_index import CompileSummaryIndex, SummaryIndexEntry

            entries: List[SummaryIndexEntry] = []
            summaries: List[str] = []

            for file_path, entry in manifest.files.items():
                if entry.summary:
                    entries.append(SummaryIndexEntry(
                        file_path=file_path,
                        summary=entry.summary,
                    ))
                    summaries.append(entry.summary)

            if not entries:
                return

            # Tokenize summaries + compute TF (always available)
            tokenizer = TokenizerUtil()
            for idx, entry in enumerate(entries):
                tokens = tokenizer.segment(entry.summary)
                entry.tokens = tokens
                entry.token_freqs = {}
                for t in tokens:
                    entry.token_freqs[t] = entry.token_freqs.get(t, 0) + 1

            # Compute embeddings (optional — requires EmbeddingUtil)
            try:
                from sirchmunk.utils.embedding_util import EmbeddingUtil
                embedding_util = EmbeddingUtil()
                embedding_util.start_loading()
                # Wait up to 60 seconds for model load
                await embedding_util._ensure_model_async(timeout=60)

                if embedding_util.is_ready():
                    embeddings = await embedding_util.embed(summaries)
                    for i, emb in enumerate(embeddings):
                        entries[i].embedding = emb
                    await self._log.info(
                        f"Summary index: computed embeddings for {len(entries)} entries"
                    )
            except Exception as emb_exc:
                await self._log.warning(
                    f"Summary index: embedding computation skipped: {emb_exc}"
                )

            index = CompileSummaryIndex(entries)
            index.save(self._compile_dir / "summary_index.json")

        except Exception as exc:
            await self._log.warning(f"Failed to build summary index: {exc}")

    # ------------------------------------------------------------------ #
    #  Manifest I/O                                                       #
    # ------------------------------------------------------------------ #

    def _load_manifest(self) -> CompileManifest:
        if self._manifest_path.exists():
            try:
                return CompileManifest.from_json(
                    self._manifest_path.read_text(encoding="utf-8")
                )
            except Exception:
                pass
        return CompileManifest()

    def _save_manifest(self, manifest: CompileManifest) -> None:
        self._manifest_path.write_text(manifest.to_json(), encoding="utf-8")

    # ------------------------------------------------------------------ #
    #  Document catalog for search-time routing                           #
    # ------------------------------------------------------------------ #

    def _build_document_catalog(self, manifest: CompileManifest) -> None:
        """Generate a lightweight catalog mapping files to their tree root summaries.

        The catalog is consumed by FAST search to fuse query analysis with
        LLM-driven document routing in a single prompt.  Each entry carries
        the filename and a truncated root summary (<= _MANIFEST_SUMMARY_MAX_LEN chars).

        Summary is sourced from the manifest (populated during Phase 2 compile),
        with a tree-root fallback for backward compatibility.
        """
        tree_cache = self._compile_dir / "trees"
        entries: List[Dict[str, str]] = []

        for file_path, entry in manifest.files.items():
            summary = entry.summary  # Primary: manifest-persisted summary

            # Fallback: read from tree root if manifest summary is empty
            if not summary and entry.has_tree and tree_cache.exists():
                tree_file = tree_cache / f"{entry.file_hash}.json"
                if tree_file.exists():
                    try:
                        tree = DocumentTree.from_json(
                            tree_file.read_text(encoding="utf-8"),
                        )
                        if tree.root and tree.root.summary:
                            summary = tree.root.summary[:_MANIFEST_SUMMARY_MAX_LEN]
                    except Exception:
                        pass

            entries.append({
                "path": file_path,
                "name": Path(file_path).name,
                "summary": summary,
            })

        catalog_path = self._compile_dir / "document_catalog.json"
        catalog_path.write_text(
            json.dumps(entries, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
